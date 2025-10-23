
from pptx import Presentation       #type: ignore
from pptx.util import Inches, Pt    #type: ignore
from pptx.enum.text import PP_ALIGN #type: ignore


def create_pptx() : 
    # --- Crée une nouvelle présentation ---
    prs = Presentation()

    # --- Ajoute une diapositive de titre ---
    slide_layout = prs.slide_layouts[0]  # 0 = mise en page "Titre + sous-titre"
    slide = prs.slides.add_slide(slide_layout)

    # --- Définit le titre ---
    title = slide.shapes.title
    title.text = "First play with images"

    # --- Optionnel : sous-titre ---
    subtitle = slide.placeholders[1]
    subtitle.text = """Learning to manipulate Images with Python 🐍
    Get familiar with OpenCV
    """

    # --- Enregistre le fichier ---
    prs.save("first_play_with_images.pptx")

    slide = prs.slides.add_slide(slide_layout)

    print("✅ PowerPoint créé : first_play_with_images.pptx")


    import numpy as np
import cv2  # type: ignore
from pptx import Presentation       # type: ignore
from pptx.util import Inches, Pt    # type: ignore
from pptx.enum.text import PP_ALIGN # type: ignore

# --- Lecture des images ---
graffiti_img = cv2.imread("Yann\\lfi-01\\graffiti.png")
print("Shape of the normal image :", graffiti_img.shape)

graffiti_img_gray = cv2.imread("Yann\\lfi-01\\graffiti.png", cv2.IMREAD_GRAYSCALE)
print("Shape of the gray image :", graffiti_img_gray.shape)

# --- Convertir la version grise en RVB pour affichage ---
graffiti_img_gray_rvb = cv2.cvtColor(graffiti_img_gray, cv2.COLOR_GRAY2BGR)

# --- Créer une image concaténée pour affichage ---
graffiti_sbs = np.concatenate((graffiti_img, graffiti_img_gray_rvb), axis=1)
cv2.imwrite("Yann\\lfi-01\\graffiti_concat.png", graffiti_sbs)

cv2.imshow("Graffiti normal", graffiti_img)
cv2.imshow("Graffiti gray", graffiti_img_gray_rvb)
cv2.imshow("Graffiti sbs", graffiti_sbs)
cv2.waitKey(2000)
cv2.destroyAllWindows()

# --- Création du PowerPoint ---
prs = Presentation()

# --- 1️⃣ Diapo de titre ---
slide_layout = prs.slide_layouts[0]  # "Titre + sous-titre"
slide = prs.slides.add_slide(slide_layout)

title = slide.shapes.title
title.text = "First play with images"

subtitle = slide.placeholders[1]
subtitle.text = """Learning to manipulate images with Python 🐍
Get familiar with OpenCV and PowerPoint automation"""

# --- 2️⃣ Nouvelle diapositive pour les images ---
blank_slide_layout = prs.slide_layouts[6]  # diapo vide
slide2 = prs.slides.add_slide(blank_slide_layout)

# --- Ajout des deux images ---
left_img = Inches(1)
top_img = Inches(1)
height_img = Inches(3)

# Image couleur
img1_path = "Yann\\lfi-01\\graffiti.png"
pic1 = slide2.shapes.add_picture(img1_path, left_img, top_img, height=height_img)

# Image gris
img2_path = "Yann\\lfi-01\\graffiti_gray.png"  # ou graffiti_img_gray si tu veux une seule image
left_img2 = Inches(5)
pic2 = slide2.shapes.add_picture(img2_path, left_img2, top_img, height=height_img)

# --- Ajout de légendes ---
# Légende 1
txBox1 = slide2.shapes.add_textbox(left_img, Inches(4.3), Inches(3), Inches(3))
tf1 = txBox1.text_frame
tf1.text = "Original color image"
tf1.paragraphs[0].alignment = PP_ALIGN.CENTER

# Légende 2
txBox2 = slide2.shapes.add_textbox(left_img2, Inches(4.3), Inches(3), Inches(3))
tf2 = txBox2.text_frame
tf2.text = "Converted grayscale image"
tf2.paragraphs[0].alignment = PP_ALIGN.CENTER

# --- Enregistrer le PowerPoint ---
prs.save("first_play_with_images.pptx")
print("✅ PowerPoint créé : first_play_with_images.pptx")

